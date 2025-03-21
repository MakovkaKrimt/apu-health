import os
from typing import List, Tuple
from dataclasses import dataclass
from pptx.util import Mm, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
from pptx.enum.text import PP_ALIGN


class PointCreator:
    def __init__(self, slide):
        self.slide = slide

    def add_site(self, point_emu: Tuple[float, float], mso_shape_type:MSO_SHAPE,shape_width_emu:float,shape_height_emu:float):
        try:
            left = point_emu[0] - shape_width_emu / 2
            top = point_emu[1] - shape_height_emu/2
            point = self.slide.shapes.add_shape(mso_shape_type, left, top, shape_width_emu, shape_height_emu)

            point.fill.solid()
            point.fill.fore_color.rgb = RGBColor(255, 180, 8)
            point.line.color.rgb = RGBColor(0, 0, 8)
            point.line.width = Pt(0.5)

            # point.rotation = -45

            return point
        except Exception as e:
            print(f"Ошибка при добавлении точки на слайд: {e}")
            return None

    def add_policlinic(self, point_emu: Tuple[float, float], mso_shape_type:MSO_SHAPE,shape_width_emu:float,shape_height_emu:float):
        try:
            left = point_emu[0] - shape_width_emu / 2
            top = point_emu[1] - shape_height_emu/2
            point = self.slide.shapes.add_shape(mso_shape_type, left, top, shape_width_emu, shape_height_emu)

            point.fill.solid()
            point.fill.fore_color.rgb = RGBColor(0, 0, 255) 

            return point
        except Exception as e:
            print(f"Ошибка при добавлении точки на слайд: {e}")
            return None
