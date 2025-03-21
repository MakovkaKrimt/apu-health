
import os
from typing import List, Tuple
from dataclasses import dataclass
from pptx.util import Mm, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
from pptx.enum.text import PP_ALIGN

from app.models.project_site_model import ProjectSiteModel


class TextboxCreator:
    def __init__(self, slide, font_name='Century Gothic', font_color = [0,0,0]):
        self.slide = slide
        self.font_name = font_name
        self.font_color = font_color

    def add_textbox(self, point_data: ProjectSiteModel, point_emu: Tuple[float, float], textbox_width: float, textbox_height: float, top:float ,left:float):
        try:
            textbox = self.slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, textbox_width, textbox_height)

            textbox.fill.solid()
            textbox.fill.fore_color.rgb = RGBColor(255, 255, 255)
            textbox.line.color.rgb = RGBColor(255, 179, 8)
            textbox.line.width = Pt(1.5)

            text_frame = textbox.text_frame
            text_frame.word_wrap = True
            text_frame.auto_size = True

            if point_data.address:
                text_frame.text = point_data.address
                paragraph = text_frame.paragraphs[0]
                paragraph.alignment = PP_ALIGN.CENTER
                paragraph.font.name = self.font_name
                paragraph.font.size = Pt(14)
                paragraph.font.color.rgb = RGBColor(*self.font_color)
                paragraph.font.italic = False
                paragraph.font.bold = True

            # if point_data.type:
            #     paragraph = text_frame.add_paragraph()
            #     paragraph.text = f"{point_data.type}  поликлиника на "
            #     paragraph.alignment = PP_ALIGN.CENTER
            #     paragraph.font.name = self.font_name
            #     paragraph.font.size = Pt(12)
            #     paragraph.font.color.rgb = RGBColor(*self.font_color)
            #     paragraph.font.italic = True
            #     paragraph.font.bold = False

            # if point_data.revenue_2023:
            #     paragraph = text_frame.add_paragraph()
            #     run1 = paragraph.add_run()
            #     run1.text = "Выручка - "
            #     run1.font.bold = False
            #     run1.font.italic = False
            #     run1.font.name = "TT Moscow Economy"
            #     run1.font.size = Pt(11.2)
            #     run1.font.color.rgb = font_color

            #     run2 = paragraph.add_run()
            #     run2.text = point_data.revenue_2023
            #     run2.font.bold = True
            #     run2.font.italic = False
            #     run2.font.name = "TT Moscow Economy"
            #     run2.font.size = Pt(11.2)
            #     run2.font.color.rgb = font_color

            #     run3 = paragraph.add_run()
            #     run3.text = " млн руб."
            #     run3.font.bold = True
            #     run3.font.italic = False
            #     run3.font.name = "TT Moscow Economy"
            #     run3.font.size = Pt(11.2)
            #     run3.font.color.rgb = font_color

            # if point_data.number_2023:
            #     paragraph = text_frame.add_paragraph()
            #     run1 = paragraph.add_run()
            #     run1.text = "Численность - "
            #     run1.font.bold = False
            #     run1.font.italic = False
            #     run1.font.name = "TT Moscow Economy"
            #     run1.font.size = Pt(11.2)
            #     run1.font.color.rgb = font_color

            #     run2 = paragraph.add_run()
            #     run2.text = point_data.number_2023
            #     run2.font.bold = True
            #     run2.font.italic = False
            #     run2.font.name = "TT Moscow Economy"
            #     run2.font.size = Pt(11.2)
            #     run2.font.color.rgb = font_color

            # if point_data.salary_2023:
            #     paragraph = text_frame.add_paragraph()
            #     run1 = paragraph.add_run()
            #     run1.text = "Средняя з/п - "
            #     run1.font.bold = False
            #     run1.font.italic = False
            #     run1.font.name = "TT Moscow Economy"
            #     run1.font.size = Pt(11.2)
            #     run1.font.color.rgb = font_color

            #     run2 = paragraph.add_run()
            #     run2.text = str(point_data.salary_2023)
            #     run2.font.bold = True
            #     run2.font.italic = False
            #     run2.font.name = "TT Moscow Economy"
            #     run2.font.size = Pt(11.2)
            #     run2.font.color.rgb = font_color

            #     run3 = paragraph.add_run()
            #     run3.text = " тыc. руб."
            #     run3.font.bold = True
            #     run3.font.italic = False
            #     run3.font.name = "TT Moscow Economy"
            #     run3.font.size = Pt(11.2)
            #     run3.font.color.rgb = font_color

            # if point_data.taxes_2023:
            #     paragraph = text_frame.add_paragraph()
            #     run1 = paragraph.add_run()
            #     run1.text = "Налоги - "
            #     run1.font.bold = False
            #     run1.font.italic = False
            #     run1.font.name = "TT Moscow Economy"
            #     run1.font.size = Pt(11.2)
            #     run1.font.color.rgb = font_color

            #     run2 = paragraph.add_run()
            #     run2.text = str(point_data.taxes_2023)
            #     run2.font.bold = True
            #     run2.font.italic = False
            #     run2.font.name = "TT Moscow Economy"
            #     run2.font.size = Pt(11.2)
            #     run2.font.color.rgb = font_color

            #     run3 = paragraph.add_run()
            #     run3.text = " млн руб."
            #     run3.font.bold = True
            #     run3.font.italic = False
            #     run3.font.name = "TT Moscow Economy"
            #     run3.font.size = Pt(11.2)
            #     run3.font.color.rgb = font_color

            return textbox
        except Exception as e:
            print(f"Ошибка при добавлении текстбокса на слайд: {e}")
            return None
