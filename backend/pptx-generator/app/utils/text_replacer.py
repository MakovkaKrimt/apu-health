from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt
from typing import Optional

class TextReplacer:
    def __init__(self, slide):
        self.slide = slide

    def find_shape_by_name(self, shapes, shape_name: str):
        for shape in shapes:
            if shape.name == shape_name:
                return shape
            
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                result = self.find_shape_by_name(shape.shapes, shape_name)
                if result:
                    return result
        
        return None

    def replace_text(self, shape_name: str, new_text: str, font_name: str = 'Century Gothic', font_size: Pt = Pt(14), bold: bool = True, italic: bool = False):
        shape = self.find_shape_by_name(self.slide.shapes, shape_name)
        
        if shape:
            # print(f"Фигура найдена: {shape.name}")
            
            text_frame = shape.text_frame
            text_frame.text = str(new_text)
            paragraph = text_frame.paragraphs[0]
            paragraph.font.name = font_name
            paragraph.font.size = font_size
            paragraph.font.bold = bold
            paragraph.font.italic = italic
            
            # print("Текст успешно изменен.")
            return True
        else:
            # print(f"Фигура с именем '{shape_name}' не найдена.")
            return False
