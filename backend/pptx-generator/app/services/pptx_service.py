from pptx import Presentation
from pptx.util import Cm
import os
import time
from dataclasses import dataclass
from typing import List, Dict, Tuple
from app.models.analysis_model import AnalysisModel
from app.models.policlinics_model import PoliclinicsModel
from app.models.population_model import PopulationDTO
from app.models.project_site_model import ProjectSiteModel
from app.utils import PointManager,TextManager,TextReplacer,CalloutManager


PAGE_HEIGHT = 297.0

class PptxService:

    async def generate_pptx(self, template_path:str,background_image_path:str, output_path: str, analysis_model:AnalysisModel,policlinics_model:PoliclinicsModel,project_site_model:ProjectSiteModel):

        prs = Presentation(template_path)
        slide = prs.slides[0]

        # print('lhlhl',analysis_model)

        text_manager = TextManager(slide)
        point_manager = PointManager(slide,PAGE_HEIGHT)
        callout_manager = CalloutManager(slide,PAGE_HEIGHT)

        slide.shapes.add_picture(background_image_path, left=Cm(2.66), top=Cm(6.64), height=Cm(16.65))

        text_manager.update_analysis_placeholders(analysis_model)

        point_manager.add_policlinics(policlinics_model)

        callout_manager.add_callout(project_site_model)

        # print(project_site_model)


        # prs.save("result.pptx")
        prs.save(output_path)